import os
import re
import csv
import shutil
import smtplib
import zipfile
import xml.etree.ElementTree as ET
from datetime import datetime
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders

import openpyxl
from openpyxl import Workbook, load_workbook
from openpyxl.utils import get_column_letter
import pandas as pd
from pptx import Presentation
from .utils import get_platform, ensure_file_exists, require_windows, get_win32_module


EXCEL_DATE_FORMATS = [
    "%Y-%m-%d", "%d/%m/%Y", "%m/%d/%Y",
    "%Y/%m/%d", "%d-%m-%Y", "%m-%d-%Y",
]

def merge_excel_files(input_paths, output_path):
    target = Workbook()
    target.remove(target.active)
    for path in input_paths:
        ensure_file_exists(path)
        wb = load_workbook(path, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            new_ws = target.create_sheet(title=sheet_name)
            for row in ws.iter_rows(min_row=1, max_row=ws.max_row,
                                    max_col=ws.max_column, values_only=False):
                for cell in row:
                    new_cell = new_ws.cell(row=cell.row, column=cell.column,
                                           value=cell.value)
                    if cell.has_style:
                        new_cell.font = cell.font
                        new_cell.fill = cell.fill
                        new_cell.number_format = cell.number_format
                        new_cell.alignment = cell.alignment
        wb.close()
    target.save(output_path)
    return output_path


def split_excel_sheet(input_path, column, output_dir, sheet_name=None):
    ensure_file_exists(input_path)
    os.makedirs(output_dir, exist_ok=True)
    df = pd.read_excel(input_path, sheet_name=sheet_name)
    if column not in df.columns:
        raise ValueError(f"Column '{column}' not found in data")
    for value, group in df.groupby(column):
        safe_name = re.sub(r'[^\w\-_ ]', '_', str(value))[:60]
        out_path = os.path.join(output_dir, f"{safe_name}.xlsx")
        group.to_excel(out_path, index=False, engine="openpyxl")
    return output_dir


def compare_excel_files(file1, file2, sheet_name=None, output_path=None):
    ensure_file_exists(file1)
    ensure_file_exists(file2)
    df1 = pd.read_excel(file1, sheet_name=sheet_name)
    df2 = pd.read_excel(file2, sheet_name=sheet_name)
    df1 = df1.fillna("").astype(str)
    df2 = df2.fillna("").astype(str)
    if df1.shape != df2.shape:
        result = f"Shape mismatch: {df1.shape} vs {df2.shape}"
        if output_path:
            wb = Workbook()
            ws = wb.active
            ws.title = "Summary"
            ws.append(["File 1 Shape", str(df1.shape)])
            ws.append(["File 2 Shape", str(df2.shape)])
            ws.append(["Result", result])
            wb.save(output_path)
        return result
    diff_mask = df1 != df2
    if not diff_mask.any().any():
        result = "Files are identical"
        if output_path:
            wb = Workbook()
            ws = wb.active
            ws.title = "Summary"
            ws.append(["Result", result])
            wb.save(output_path)
        return result
    if output_path:
        wb = Workbook()
        ws = wb.active
        ws.title = "Differences"
        ws.append(["Row", "Column", "File 1 Value", "File 2 Value"])
        for r in range(len(df1)):
            for c in range(len(df1.columns)):
                if diff_mask.iloc[r, c]:
                    ws.append([
                        int(r + 2),
                        df1.columns[c],
                        df1.iloc[r, c],
                        df2.iloc[r, c],
                    ])
        wb.save(output_path)
    return f"Found {diff_mask.sum().sum()} difference(s)"


def search_excel(input_path, value, sheet_name=None):
    ensure_file_exists(input_path)
    wb = load_workbook(input_path, data_only=True)
    results = []
    search_value = str(value).lower()
    sheets = [sheet_name] if sheet_name else wb.sheetnames
    for name in sheets:
        if name not in wb:
            continue
        ws = wb[name]
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row,
                                max_col=ws.max_column, values_only=False):
            for cell in row:
                if cell.value is not None and search_value in str(cell.value).lower():
                    results.append({
                        "sheet": name,
                        "row": cell.row,
                        "column": get_column_letter(cell.column),
                        "value": cell.value,
                    })
    wb.close()
    return results


def clean_excel(input_path, output_path, remove_duplicates=True,
                remove_empty_rows=True, fix_dates=True):
    ensure_file_exists(input_path)
    wb = load_workbook(input_path)
    for ws in wb.worksheets:
        if remove_empty_rows:
            rows_to_delete = []
            for row_idx in range(ws.max_row, 0, -1):
                if all(ws.cell(row=row_idx, column=c).value is None
                       for c in range(1, ws.max_column + 1)):
                    rows_to_delete.append(row_idx)
            for r in rows_to_delete:
                ws.delete_rows(r)
        if remove_duplicates:
            seen = set()
            rows_to_delete = []
            for row_idx in range(1, ws.max_row + 1):
                row_vals = tuple(
                    ws.cell(row=row_idx, column=c).value
                    for c in range(1, ws.max_column + 1)
                )
                if row_vals in seen:
                    rows_to_delete.append(row_idx)
                else:
                    seen.add(row_vals)
            for r in reversed(rows_to_delete):
                ws.delete_rows(r)
        if fix_dates:
            for row in ws.iter_rows(min_row=1, max_row=ws.max_row,
                                    max_col=ws.max_column):
                for cell in row:
                    if isinstance(cell.value, str):
                        trimmed = cell.value.strip()
                        for fmt in EXCEL_DATE_FORMATS:
                            try:
                                dt = datetime.strptime(trimmed, fmt)
                                cell.value = dt
                                break
                            except ValueError:
                                continue
    wb.save(output_path)
    return output_path


def _get_docx_text(docx_path):
    ensure_file_exists(docx_path)
    text_parts = []
    with zipfile.ZipFile(docx_path) as z:
        with z.open("word/document.xml") as f:
            tree = ET.parse(f)
            root = tree.getroot()
            for t in root.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                if t.text:
                    text_parts.append(t.text)
    return "".join(text_parts)


def _replace_docx_placeholders(docx_path, replacements, output_path):
    ensure_file_exists(docx_path)
    shutil.copy2(docx_path, output_path)
    with zipfile.ZipFile(output_path, "r") as zin:
        doc_xml = zin.read("word/document.xml").decode("utf-8")
        for key, value in replacements.items():
            placeholder = "{{" + key + "}}"
            doc_xml = doc_xml.replace(placeholder, str(value))
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zout:
        with zipfile.ZipFile(output_path, "r") as zin:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "word/document.xml":
                    data = doc_xml.encode("utf-8")
                zout.writestr(item, data)


def merge_word_docx(template_path, data_path, output_dir):
    ensure_file_exists(template_path)
    ensure_file_exists(data_path)
    os.makedirs(output_dir, exist_ok=True)
    df = pd.read_excel(data_path)
    columns = df.columns.tolist()
    base_name = os.path.splitext(os.path.basename(template_path))[0]
    for idx, row in df.iterrows():
        replacements = {col: str(row[col]) for col in columns}
        out_path = os.path.join(output_dir, f"{base_name}_{idx + 1}.docx")
        _replace_docx_placeholders(template_path, replacements, out_path)
    return output_dir


def convert_docx(input_path, output_path):
    ensure_file_exists(input_path)
    ext = os.path.splitext(output_path)[1].lower()
    if ext == ".txt":
        text = _get_docx_text(input_path)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
    elif ext == ".pdf":
        raise RuntimeError(
            "DOCX to PDF conversion requires LibreOffice or Word on the system. "
            "Use .txt output for cross-platform text extraction."
        )
    else:
        raise ValueError(f"Unsupported output format: {ext}")
    return output_path


def search_emails(subject=None, sender=None, since=None, until=None,
                  folder="Inbox", max_results=50):
    require_windows("Outlook search")
    win32com = get_win32_module("win32com.client")
    outlook = win32com.Dispatch("Outlook.Application").GetNamespace("MAPI")
    folder_obj = outlook.Folders.Item(1).Folders.Item(folder)
    messages = folder_obj.Items
    if since:
        messages = messages.Restrict(f"[ReceivedTime] >= '{since}'")
    if until:
        messages = messages.Restrict(f"[ReceivedTime] <= '{until}'")
    results = []
    count = 0
    for msg in messages:
        if count >= max_results:
            break
        try:
            msg_subject = msg.Subject or ""
            msg_sender = msg.SenderName or msg.SenderEmailAddress or ""
            if subject and subject.lower() not in msg_subject.lower():
                continue
            if sender and sender.lower() not in msg_sender.lower():
                continue
            results.append({
                "subject": msg_subject,
                "sender": msg_sender,
                "received": str(msg.ReceivedTime),
                "body_preview": (msg.Body or "")[:200],
            })
            count += 1
        except Exception:
            continue
    return results


def send_email(smtp_server, smtp_port, username, password,
               to_addrs, subject, body, attachments=None,
               use_tls=True):
    msg = MIMEMultipart()
    msg["From"] = username
    msg["To"] = ", ".join(to_addrs) if isinstance(to_addrs, list) else to_addrs
    msg["Subject"] = subject
    msg.attach(MIMEText(body, "plain", "utf-8"))
    if attachments:
        for filepath in attachments:
            ensure_file_exists(filepath)
            with open(filepath, "rb") as f:
                part = MIMEBase("application", "octet-stream")
                part.set_payload(f.read())
                encoders.encode_base64(part)
                filename = os.path.basename(filepath)
                part.add_header(
                    "Content-Disposition",
                    f"attachment; filename=\"{filename}\""
                )
                msg.attach(part)
    if use_tls:
        server = smtplib.SMTP(smtp_server, smtp_port)
        server.starttls()
    else:
        server = smtplib.SMTP_SSL(smtp_server, smtp_port)
    server.login(username, password)
    to_list = [to_addrs] if isinstance(to_addrs, str) else to_addrs
    server.sendmail(username, to_list, msg.as_string())
    server.quit()
    return True


def export_emails(output_path, subject=None, sender=None,
                  since=None, until=None, folder="Inbox", max_results=200):
    emails = search_emails(subject, sender, since, until, folder, max_results)
    ext = os.path.splitext(output_path)[1].lower()
    if ext in (".csv",):
        with open(output_path, "w", newline="", encoding="utf-8") as f:
            if emails:
                writer = csv.DictWriter(f, fieldnames=emails[0].keys())
                writer.writeheader()
                writer.writerows(emails)
    elif ext in (".xlsx",):
        df = pd.DataFrame(emails)
        df.to_excel(output_path, index=False, engine="openpyxl")
    else:
        raise ValueError(f"Unsupported export format: {ext}")
    return output_path


def batch_ppt(template_path, data_path, output_dir):
    ensure_file_exists(template_path)
    ensure_file_exists(data_path)
    os.makedirs(output_dir, exist_ok=True)
    df = pd.read_excel(data_path)
    base_name = os.path.splitext(os.path.basename(template_path))[0]
    for idx, row in df.iterrows():
        prs = Presentation(template_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            for col in df.columns:
                                placeholder = "{{" + col + "}}"
                                if placeholder in text:
                                    run.text = text.replace(
                                        placeholder, str(row[col])
                                    )
        out_path = os.path.join(output_dir, f"{base_name}_{idx + 1}.pptx")
        prs.save(out_path)
    return output_dir


def export_ppt(input_path, output_dir, format_type="png"):
    ensure_file_exists(input_path)
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(input_path)
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    exported = []
    for i, slide in enumerate(prs.slides, start=1):
        if format_type in ("png", "jpg", "jpeg"):
            ext = "png" if format_type == "png" else "jpg"
            out_path = os.path.join(output_dir, f"{base_name}_slide_{i}.{ext}")
            try:
                image = slide.shapes[0].image
                with open(out_path, "wb") as f:
                    f.write(image.blob)
            except (IndexError, AttributeError):
                with open(out_path, "w") as f:
                    f.write(f"Slide {i} - {slide.name}\n")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            f.write(shape.text_frame.text + "\n")
            exported.append(out_path)
        elif format_type == "pdf":
            out_path = os.path.join(output_dir, f"{base_name}_slide_{i}.txt")
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(f"Slide {i}\n")
                f.write("=" * 40 + "\n")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        f.write(shape.text_frame.text + "\n\n")
            exported.append(out_path)
    return exported
